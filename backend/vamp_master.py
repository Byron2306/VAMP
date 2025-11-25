#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
vamp_master.py — Monthly Scan → Parse → Score (NWU Brain–aligned)

This script orchestrates evidence ingestion and deterministic scoring for the
NWU VAMP performance management workflow. It preserves your existing flow
(file-system scan, robust text extraction, CSV export) but **replaces all
legacy routing/scoring** with the single source of truth:

    nwu_brain/brain_manifest.json  +  nwu_brain.scoring.NWUScorer

Key changes (vs older versions)
-------------------------------
1) Manifest-first loading
   • NO references to legacy input/*.json or "*updated (3).json".
   • The scorer is created from brain_manifest.json and used everywhere.

2) Deterministic scoring (single source)
   • For each artefact we call `SCORER.compute(…)`.
   • CSV export uses `SCORER.to_csv_row(…)` (CSV v2). You can optionally
     keep a few legacy columns for downstream compatibility.

3) Deep Read Mode by default
   • Extraction tries to read **full text** (PDF, DOCX, XLSX, PPTX, TXT/CSV/MD).
   • If extraction fails, we still record filename so the pipeline doesn’t break.

4) Non-destructive
   • Folders and filenames remain the same.
   • Connectors are outside this script; this only operates on a local root.

Invocation
----------
$ python vamp_master.py --root ./VAMP --year 2025 --month 8
# or interactive if args omitted

Optional environment variables
------------------------------
SCAN_START=YYYY-MM-DDTHH:MM:SS
SCAN_END=YYYY-MM-DDTHH:MM:SS
    If provided, these bound the scan window regardless of --year/--month.

VAMP_EMPLOYEE_NAME="Your Name"
VAMP_EMPLOYEE_RANK="Lecturer"
    Used only for informational report footers.

Outputs
-------
• <root>/_out/audit.csv        (CSV v2 rows from NWUScorer)
• <root>/_out/scan_report.md   (brief scan summary)

After all months are done, you can run vamp_runner.py to aggregate.
"""

from __future__ import annotations

import argparse
import csv
import datetime
import json
import os
import shutil
import sys
import zipfile
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional, Tuple

# -------------------------
# Paths & Brain scaffolding
# -------------------------

from . import BRAIN_DATA_DIR
from .agent_app.app_state import agent_state
from .vamp_agent_bridge import submit_evidence_from_vamp
from .nwu_brain.scoring import NWUScorer

# -------------------------
# Paths & Brain scaffolding
# -------------------------

ROOT_DIR = Path(__file__).resolve().parent
MANIFEST_PATH = BRAIN_DATA_DIR / "brain_manifest.json"

if not MANIFEST_PATH.is_file():
    raise FileNotFoundError(f"Brain manifest not found: {MANIFEST_PATH}")

SCORER = NWUScorer(str(MANIFEST_PATH))

# -------------------------
# Utilities & helpers
# -------------------------

def say(msg: str) -> None:
    print(msg, flush=True)

def ensure_dir(p: Path) -> Path:
    p.mkdir(parents=True, exist_ok=True)
    return p

def month_bounds(year: int, month: int) -> Tuple[datetime.datetime, datetime.datetime]:
    """Return naive UTC-like bounds for the calendar month."""
    start = datetime.datetime(year, month, 1)
    end = datetime.datetime(year + (month // 12), (month % 12) + 1, 1)
    return start, end

def env_scan_window() -> Tuple[Optional[datetime.datetime], Optional[datetime.datetime]]:
    """If SCAN_START/SCAN_END are set, parse them and return window; else (None,None)."""
    s = os.getenv("SCAN_START", "").strip()
    e = os.getenv("SCAN_END", "").strip()
    if not s or not e:
        return None, None
    try:
        start = datetime.datetime.fromisoformat(s)
        end = datetime.datetime.fromisoformat(e)
        return start, end
    except Exception:
        say("Warning: could not parse SCAN_START/SCAN_END; ignoring.")
        return None, None

def in_window(ts: float, start: Optional[datetime.datetime], end: Optional[datetime.datetime]) -> bool:
    """Check file mtime (POSIX seconds) falls inside window; permissive if not configured."""
    if not start or not end:
        return True
    mt = datetime.datetime.fromtimestamp(ts)
    return start <= mt < end

def sha1_file(path: Path) -> str:
    import hashlib
    h = hashlib.sha1()
    with path.open("rb") as fp:
        for chunk in iter(lambda: fp.read(65536), b""):
            h.update(chunk)
    return h.hexdigest()

def guess_relpath(root: Path, p: Path) -> str:
    try:
        return str(p.relative_to(root))
    except Exception:
        return p.name

# -------------------------
# Text extraction (Deep Read)
# -------------------------

def _bytes_decode_guess(raw: bytes) -> str:
    try:
        import chardet  # type: ignore
        enc = chardet.detect(raw).get("encoding") or "utf-8"
        return raw.decode(enc, errors="ignore")
    except Exception:
        return raw.decode("utf-8", errors="ignore")

def txt_from_pdf(path: Path) -> str:
    """Extract text & table-ish content; robust to mildly corrupted PDFs."""
    text = ""
    # Attempt to validate/repair
    try:
        import pikepdf  # type: ignore
        with pikepdf.open(str(path)):
            pass
    except Exception:
        pass
    # pdfminer
    try:
        from pdfminer.high_level import extract_text  # type: ignore
        text = extract_text(str(path)) or ""
    except Exception:
        text = ""
    # pdfplumber tables
    try:
        import pdfplumber  # type: ignore
        with pdfplumber.open(str(path)) as pdf:
            rows: List[str] = []
            for pg in pdf.pages:
                try:
                    tables = pg.extract_tables() or []
                    for tbl in tables:
                        rows.extend([" ".join([c or "" for c in r]) for r in tbl])
                except Exception:
                    pass
            if rows:
                text += "\n" + "\n".join(rows)
    except Exception:
        pass
    return text

def txt_from_docx(path: Path) -> str:
    try:
        import docx  # type: ignore
        d = docx.Document(str(path))
        return "\n".join(p.text for p in d.paragraphs)
    except Exception:
        return ""

def txt_from_xlsx(path: Path) -> str:
    try:
        import openpyxl  # type: ignore
        wb = openpyxl.load_workbook(str(path), data_only=True, read_only=True)
        out: List[str] = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                out.append(" ".join("" if c is None else str(c) for c in row))
        return "\n".join(out)
    except Exception:
        return ""

def txt_from_pptx(path: Path) -> str:
    try:
        from pptx import Presentation  # type: ignore
        prs = Presentation(str(path))
        out: List[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    out.append(shape.text or "")
        return "\n".join(out)
    except Exception:
        return ""

def extract_text_for(path: Path, size_limit: int = 200_000) -> str:
    """Deep Read extraction; on failure returns filename only."""
    name = path.name.lower()
    try:
        if name.endswith(".pdf"):
            s = txt_from_pdf(path)
        elif name.endswith(".docx"):
            s = txt_from_docx(path)
        elif name.endswith(".xlsx"):
            s = txt_from_xlsx(path)
        elif name.endswith(".pptx"):
            s = txt_from_pptx(path)
        elif name.endswith((".txt", ".md", ".csv", ".log")):
            s = _bytes_decode_guess(path.read_bytes())
        elif name.endswith(".zip"):
            try:
                with zipfile.ZipFile(str(path), "r") as zf:
                    inner = zf.namelist()[:80]
                    s = "ZIP " + " | ".join(inner)
            except Exception:
                s = name
        else:
            s = name
    except Exception:
        s = name
    s = s or name
    if len(s) > size_limit:
        s = s[:size_limit]
    return s

# -------------------------
# Ingestion
# -------------------------

SKIP_DIRS = {"_out", "_final", "_logs", ".git", "__pycache__"}

@dataclass
class Artefact:
    path: Path
    relpath: str
    size: int
    mtime: float
    sha1: str

def ingest_paths(evidence_root: Path,
                 start: Optional[datetime.datetime],
                 end: Optional[datetime.datetime]) -> List[Artefact]:
    items: List[Artefact] = []
    seen: set[str] = set()
    for root, dirs, files in os.walk(evidence_root):
        # prune skip dirs in-place for performance
        dirs[:] = [d for d in dirs if d not in SKIP_DIRS]
        for fn in files:
            p = Path(root) / fn
            try:
                st = p.stat()
            except Exception:
                continue
            if not in_window(st.st_mtime, start, end):
                continue
            try:
                h = sha1_file(p)
            except Exception:
                h = f"ERR::{p.as_posix()}::{st.st_size}"
            if h in seen:
                continue
            seen.add(h)
            items.append(Artefact(
                path=p,
                relpath=guess_relpath(evidence_root, p),
                size=st.st_size,
                mtime=st.st_mtime,
                sha1=h,
            ))
    return items

# -------------------------
# Scoring & CSV export
# -------------------------

CSV_V2_FIELDS_ORDER = [
    # Common/legacy-ish context (kept for compatibility)
    "name", "relpath", "platform", "source",
    "size", "modified", "hash",
    # Deterministic NWU Brain fields
    "kpa", "tier", "tier_rule",
    "values_score", "values_hits",
    "policy_hits", "policy_hit_details", "must_pass_risks",
    "score", "band",
    "rationale", "actions",
]

def to_common_context_row(artefact: Artefact, item: Dict[str, Any]) -> Dict[str, Any]:
    """
    Build a compatibility context row + merge with scorer CSV v2 fields.
    The scorer row is authoritative for policy/kpa/tier/score/band/etc.
    """
    ctx = {
        "name": artefact.path.name,
        "relpath": artefact.relpath,
        "platform": item.get("platform", "LocalFS"),
        "source": item.get("source", "fs"),
        "size": artefact.size,
        "modified": datetime.datetime.fromtimestamp(artefact.mtime).isoformat(timespec="seconds"),
        "hash": artefact.sha1,
    }
    return ctx

def write_csv_v2(out_csv: Path, rows: List[Dict[str, Any]]) -> None:
    ensure_dir(out_csv.parent)
    # union of provided fields with ordered header preferring CSV_V2_FIELDS_ORDER
    all_keys = list(dict.fromkeys(CSV_V2_FIELDS_ORDER + [k for r in rows for k in r.keys()]))
    with out_csv.open("w", newline="", encoding="utf-8") as fp:
        w = csv.DictWriter(fp, fieldnames=all_keys, extrasaction="ignore")
        w.writeheader()
        for r in rows:
            w.writerow(r)

def write_scan_report_md(out_md: Path, scanned_count: int, used_window: Tuple[Optional[datetime.datetime], Optional[datetime.datetime]]) -> None:
    ensure_dir(out_md.parent)
    s, e = used_window
    with out_md.open("w", encoding="utf-8") as fp:
        fp.write("# VAMP Monthly Scan Report\n\n")
        fp.write(f"- Scanned items (deduped): **{scanned_count}**\n")
        if s and e:
            fp.write(f"- Window: **{s.isoformat()} → {e.isoformat()}**\n")
        else:
            fp.write("- Window: **(not bounded; full scan)**\n")
        emp = os.getenv("VAMP_EMPLOYEE_NAME", "").strip()
        rnk = os.getenv("VAMP_EMPLOYEE_RANK", "").strip()
        if emp or rnk:
            fp.write(f"- Employee: **{emp or 'N/A'}**  —  Rank: **{rnk or 'N/A'}**\n")
        fp.write("\nGenerated by vamp_master.py (NWU Brain–aligned).\n")

# -------------------------
# Main scanning routine
# -------------------------

def scan_and_score(evidence_root: Path,
                   year: int,
                   month: int,
                   out_dirname: str = "_out") -> Tuple[Path, Path]:
    """
    Walk the evidence_root (including any online sync folders you might have placed
    inside it), extract text, score deterministically, and export CSV v2.
    Returns: (audit_csv_path, report_md_path)
    """
    # Decide time window
    env_s, env_e = env_scan_window()
    if env_s and env_e:
        s, e = env_s, env_e
    else:
        s, e = month_bounds(year, month)

    say(f"Evidence root: {evidence_root}")
    say(f"Scan window:  {s.isoformat()} → {e.isoformat()}")

    artefacts = ingest_paths(evidence_root, s, e)
    say(f"Found {len(artefacts)} candidate files in window.")

    rows_csv: List[Dict[str, Any]] = []

    state = agent_state()

    for i, art in enumerate(artefacts, 1):
        if i % 25 == 1 or i == len(artefacts):
            say(f"  · Reading/scoring {i}/{len(artefacts)} … {art.relpath}")

        full_text = extract_text_for(art.path)

        # Build a minimal item for scorer; do NOT re-implement policy/KPA logic
        item = {
            "source": "fs",
            "platform": "LocalFS",
            "title": art.path.name,
            "path": str(art.path),
            "relpath": art.relpath,
            "size": art.size,
            "modified": datetime.datetime.fromtimestamp(art.mtime).isoformat(timespec="seconds"),
            "hash": art.sha1,
            "full_text": full_text,   # Deep Read text (critical for policy hits)
        }

        try:
            scored = SCORER.compute(item)                # canonical compute
            csv_row = SCORER.to_csv_row(scored)          # CSV v2 columns
        except Exception as ex:
            # If scoring fails, keep a degraded row so the pipeline continues
            csv_row = {
                "kpa": [],
                "tier": [],
                "tier_rule": "",
                "values_score": 0.0,
                "values_hits": [],
                "policy_hits": [],
                "policy_hit_details": [],
                "must_pass_risks": [],
                "score": 0.0,
                "band": "",
                "rationale": f"(scoring error: {ex})",
                "actions": [],
            }

        # Merge some legacy/common context fields for compatibility
        ctx = to_common_context_row(artefact=art, item=item)
        out_row = {**ctx, **csv_row}
        rows_csv.append(out_row)

        # Persist to the agent-managed evidence vault for audit & retention control
        try:
            state.record_evidence(
                {
                    "uid": out_row.get("hash") or art.sha1,
                    "source": out_row.get("platform", "LocalFS"),
                    "title": out_row.get("name", art.path.name),
                    "kpas": out_row.get("kpa", []) or [],
                    "score": float(out_row.get("score", 0.0) or 0.0),
                    "rationale": out_row.get("rationale", ""),
                    "metadata": {
                        "relpath": out_row.get("relpath"),
                        "policy_hits": out_row.get("policy_hits", []),
                        "must_pass_risks": out_row.get("must_pass_risks", []),
                    },
                }
            )
        except Exception as exc:  # pragma: no cover - non-critical path
            say(f"Warning: failed to persist evidence to vault: {exc}")

        submit_evidence_from_vamp(
            {
                "path": art.path,
                "evidence_id": out_row.get("hash") or art.sha1,
                "text": full_text,
                "kpa": out_row.get("kpa", []),
                "score": out_row.get("score", 0.0),
                "metadata": {
                    "source": "vamp_master",
                    "relpath": out_row.get("relpath"),
                },
            }
        )

    out_dir = ensure_dir(evidence_root / out_dirname)
    audit_csv = out_dir / "audit.csv"
    report_md = out_dir / "scan_report.md"

    write_csv_v2(audit_csv, rows_csv)
    write_scan_report_md(report_md, len(artefacts), (s, e))

    say(f"Wrote CSV → {audit_csv}")
    say(f"Wrote report → {report_md}")
    return audit_csv, report_md

# -------------------------
# CLI / main
# -------------------------

def parse_args(argv: Optional[List[str]] = None) -> argparse.Namespace:
    p = argparse.ArgumentParser(description="VAMP Monthly Scanner (NWU Brain–aligned)")
    p.add_argument("--root", type=str, default=str(ROOT_DIR.parent),
                   help="Evidence root folder (default: repository root)")
    p.add_argument("--year", type=int, default=datetime.datetime.now().year,
                   help="Assessment year (default: current year)")
    p.add_argument("--month", type=int, default=datetime.datetime.now().month,
                   help="Month 1-12 (default: current month)")
    p.add_argument("--out", type=str, default="_out", help="Output folder name inside root (default: _out)")
    p.add_argument("--mkdir", action="store_true", help="Create root if it doesn't exist")
    return p.parse_args(argv)

def main(argv: Optional[List[str]] = None) -> int:
    args = parse_args(argv)

    evidence_root = Path(args.root).resolve()
    if not evidence_root.exists():
        if args.mkdir:
            ensure_dir(evidence_root)
        else:
            say(f"ERROR: root does not exist: {evidence_root}")
            return 2

    try:
        scan_and_score(evidence_root=evidence_root, year=args.year, month=args.month, out_dirname=args.out)
    except KeyboardInterrupt:
        say("\nInterrupted.")
        return 130
    except Exception as e:
        say(f"FATAL: {e}")
        return 1
    return 0

if __name__ == "__main__":
    sys.exit(main())
